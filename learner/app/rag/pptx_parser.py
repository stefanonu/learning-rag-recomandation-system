from pptx import Presentation

def parse_pptx(file_path: str) -> list[str]:
    prs = Presentation(file_path)
    chunks = []

    for slide in prs.slides:
        bullets = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    bullets.append(text)

        # Merge small bullet points into ~100-word chunks
        current_chunk = ""
        for bullet in bullets:
            if len(current_chunk.split()) + len(bullet.split()) < 80:
                current_chunk += " " + bullet
            else:
                chunks.append(current_chunk.strip())
                current_chunk = bullet
        if current_chunk:
            chunks.append(current_chunk.strip())

    return chunks
